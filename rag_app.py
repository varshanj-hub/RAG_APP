import streamlit as st
import tempfile
import shutil
from pathlib import Path
import zipfile
import re
import nltk

from pypdf import PdfReader
import docx
import pptx
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np

from langchain_chroma import Chroma
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.chains import RetrievalQA
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_groq import ChatGroq

# Streamlit Page Config
st.set_page_config("RAG Assistant", page_icon="🔍", layout="wide")

# CSS Styling for Clarity and Professionalism
st.markdown('''
    <style>
    body {
        font-family: 'Segoe UI', sans-serif;
    }
    .header {
        font-size: 30px;
        font-weight: bold;
        color: #234;
        text-align: center;
        margin-top: 1rem;
    }
    .subheader {
        font-size: 18px;
        color: #456;
        margin: 1rem 0;
    }
    .context-box {
        background-color: #f8f9fa;
        padding: 1rem;
        border-radius: 10px;
        margin-bottom: 1rem;
        border-left: 5px solid #1f77b4;
        font-family: monospace;
        color: black;
    }
    </style>
''', unsafe_allow_html=True)

st.markdown("<div class='header'>📄 RAG Q&A Assistant</div>", unsafe_allow_html=True)

# Extractors

def extract_pdf(path):
    text = ""
    with open(path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            t = page.extract_text()
            if t: text += t + "\n"
    return text

def extract_docx(path):
    doc = docx.Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_pptx(path):
    prs = pptx.Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"): text += shape.text + "\n"
    return text

def extract_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".txt": extract_txt,
}

def clean_text(text):
    text = ''.join(c for c in text if c.isprintable() or c in '\n\t ')
    text = re.sub(r'[ \t]{2,}', ' ', text)
    text = re.sub(r'\n{2,}', '\n', text)
    text = '\n'.join(line.strip() for line in text.splitlines())
    return '\n'.join(line for line in text.splitlines() if line.strip())

def semantic_chunk(text, model, percentile_threshold=95, min_chunk_size=3, max_words=1500, overlap=200):
    sentences = nltk.sent_tokenize(text)
    if not sentences: return []
    contextualized = [' '.join(sentences[max(0, i-1):i+2]) for i in range(len(sentences))]
    embeddings = model.encode(contextualized, batch_size=16, show_progress_bar=False)
    distances = [1 - cosine_similarity([embeddings[i]], [embeddings[i+1]])[0][0] for i in range(len(embeddings)-1)]
    threshold = np.percentile(distances, percentile_threshold)
    breakpoints = [i for i, dist in enumerate(distances) if dist > threshold]

    chunks, start = [], 0
    for bp in breakpoints:
        chunks.append(' '.join(sentences[start:bp+1]))
        start = bp + 1
    if start < len(sentences): chunks.append(' '.join(sentences[start:]))

    merged_chunks, buffer, count = [], "", 0
    for chunk in chunks:
        wc = len(chunk.split())
        if count + wc < min_chunk_size * 10:
            buffer += " " + chunk
            count += wc
            continue
        words = (buffer + " " + chunk).strip().split()
        i = 0
        while i < len(words):
            merged_chunks.append(' '.join(words[i:max_words+i]))
            i += max_words - overlap
        buffer, count = "", 0
    if buffer: merged_chunks.append(buffer.strip())
    return [c for c in merged_chunks if c.strip()]

# Sidebar Configuration
st.sidebar.header("Configuration")
provider = st.sidebar.selectbox("LLM Provider", ["Gemini", "Groq"])
api_key = st.sidebar.text_input(f"{provider} API Key", type="password")
model_name = st.sidebar.selectbox("Model", [
    "gemini-2.5-pro" if provider == "Gemini" else "llama3-70b-8192",
    "Other"
])
if model_name == "Other":
    model_name = st.sidebar.text_input("Enter model name")

# Advanced Settings Dropdown
with st.sidebar.expander("⚙️ Advanced Settings"):
    chunk_size = st.slider("Max Words per Chunk", min_value=100, max_value=2000, value=1500, step=100)
    chunk_overlap = st.slider("Overlap between Chunks", min_value=0, max_value=500, value=200, step=50)
    min_sentences = st.slider("Min Sentences per Chunk", min_value=1, max_value=10, value=3)

# Document Upload and Processing
st.subheader("1. Upload and Process Files")
files = st.file_uploader("Upload PDF, DOCX, PPTX, or TXT files", type=["pdf", "docx", "pptx", "txt"], accept_multiple_files=True)
zip_file = st.file_uploader("Or upload a ZIP folder containing documents", type=["zip"])

if st.button("Process Documents") and (files or zip_file):
    embedder = SentenceTransformer('all-MiniLM-L6-v2')
    chunks, metadatas, ids = [], [], []
    with tempfile.TemporaryDirectory() as tmpdir:
        file_paths = []
        if files:
            for file in files:
                path = Path(tmpdir) / file.name
                with open(path, 'wb') as f:
                    f.write(file.getbuffer())
                file_paths.append(path)

        if zip_file:
            zip_path = Path(tmpdir) / "uploaded.zip"
            with open(zip_path, 'wb') as f:
                f.write(zip_file.getbuffer())
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(tmpdir)
            for ext in EXTRACTORS:
                file_paths.extend(Path(tmpdir).rglob(f"*{ext}"))

        for i, path in enumerate(file_paths):
            extractor = EXTRACTORS.get(path.suffix.lower())
            if extractor:
                raw = extractor(path)
                text = clean_text(raw)
                parts = semantic_chunk(text, embedder, min_chunk_size=min_sentences, max_words=chunk_size, overlap=chunk_overlap)
                for j, p in enumerate(parts):
                    chunks.append(p)
                    metadatas.append({"source": path.name, "chunk_id": j+1})
                    ids.append(f"{path.name}_chunk{j+1}")

    st.session_state["vector_store"] = Chroma.from_texts(
        texts=chunks,
        embedding=HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2"),
        metadatas=metadatas,
        ids=ids,
        persist_directory="./chroma_rag_store"
    )
    st.success(f"Indexed {len(chunks)} text chunks from uploaded documents.")

# QA Section
st.subheader("2. Ask Your Documents")
if "vector_store" in st.session_state and api_key and model_name:
    query = st.text_area("Enter your question:")
    if st.button("Ask") and query:
        llm = ChatGoogleGenerativeAI(model=model_name, google_api_key=api_key) if provider == "Gemini" else ChatGroq(groq_api_key=api_key, model=model_name)
        retriever = st.session_state['vector_store'].as_retriever(search_kwargs={"k": 3})
        qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=retriever)
        answer = qa_chain.run(query)
        st.markdown("### Answer")
        st.info(answer)
        st.markdown("#### Context Chunks")
        results = retriever.get_relevant_documents(query)
        for i, doc in enumerate(results):
            st.markdown(f"<div class='context-box'><b>Chunk {i+1} — {doc.metadata.get('source')}</b><br>{doc.page_content[:600]}{'...' if len(doc.page_content) > 600 else ''}</div>", unsafe_allow_html=True)
else:
    st.info("Please upload documents and provide valid API credentials to query your data.")
